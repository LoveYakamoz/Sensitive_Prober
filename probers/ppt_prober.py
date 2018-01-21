from pptx import Presentation

from probers.prober import Prober
from utils.log import logger


class Ppt_Prober(Prober):
    """
    pdf文档检测
    """

    def __init__(self, file, sensitive_list):
        super(Ppt_Prober, self).__init__(file, sensitive_list)

    def start_probe(self):
        print("=====> %s" % self._file)

        if self._file.endswith("ppt"):
            logger.warn("not support file type, file: %s", self._file)
            return

        def read_ppt(ppt):
            prs = Presentation(ppt)
            text_runs = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_runs.append(run.text)
            return text_runs

        try:
            text_runs = read_ppt(self._file)
            for text in text_runs:
                for sensitive in self._sensitive_list:
                    if sensitive in text and sensitive not in self._result_list:
                        self._result_list.append(sensitive)
        except Exception as e:
            logger.error(e)

    def end_probe(self):
        if len(self._result_list) > 0:
            logger.info(self)
        else:
            pass
